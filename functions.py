from langchain.schema import HumanMessage, SystemMessage
from langchain.chat_models.gigachat import GigaChat
from settings import TOKEN
from langchain_community.document_loaders import PyPDFLoader
import PyPDF2
from docx import Document
import os
import uuid
from pptx import Presentation
from pptx.util import Inches
from app.helpers import generate_image, get_files_dir
import json
import os
import uuid
import requests
import re
import urllib3
from settings import TOKEN
from pptx.dml.color import RGBColor

chat = GigaChat(credentials=TOKEN, verify_ssl_certs=False, scope="GIGACHAT_API_CORP")




urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)


def image_request(context, user_request) -> str:
    url = "https://gigachat.devices.sberbank.ru/api/v1/chat/completions"

    payload = json.dumps({
        "model": "GigaChat",
        "messages": [
            {
                "role": "user",
                "content": context + "\n" + user_request
            }
        ],

        "temperature": 1,
        "top_p": 0.1,
        "n": 1,
        "stream": False,
        "max_tokens": 512,
        "repetition_penalty": 1,
        "update_interval": 0
    })
    headers = {
        'Content-Type': 'application/json',
        'Accept': 'application/json',
        'Authorization': 'Bearer ' + get_access_token()
    }

    response = requests.request("POST", url, headers=headers, data=payload, verify=False)

    print(response.text)
    token = response.json()['choices'][0]['message']['content']
    pattern = r'src="(\w{8}-\w{4}-\w{4}-\w{4}-\w{12})"'
    match = re.search(pattern, token)
    return match.group(1)


def generate_image(slide_text: str):
    context = "Ты профессиональный дизайнер изображений для слайдов. Картинка обязательно должна быть светлой. Картинка обязательно должна быть на 75% прозрачна"
    prompt = f"Создай картинку для слайда с текстом {slide_text}"
    image_id = image_request(context, prompt)
    url = f"https://gigachat.devices.sberbank.ru/api/v1/files/{image_id}/content"

    payload = {}
    headers = {
        'Accept': 'application/jpg',
        'Authorization': 'Bearer ' + get_access_token()
    }

    response = requests.request("GET", url, headers=headers, data=payload, verify=False)
    path = os.path.join(get_files_dir(), f"{image_id}.jpg")
    with open(path, "wb") as bin_file:
        bin_file.write(response.content)

    return path


def get_access_token() -> str:
    url = "https://ngw.devices.sberbank.ru:9443/api/v2/oauth"
    payload = 'scope=GIGACHAT_API_CORP'
    headers = {
        'Content-Type': 'application/x-www-form-urlencoded',
        'Accept': 'application/json',
        'RqUID': str(uuid.uuid4()),
        'Authorization': 'Basic ' + TOKEN
    }

    response = requests.request("POST", url, headers=headers, data=payload, verify=False)
    return response.json()['access_token']


def get_files_dir() -> str:
    root_dir = os.getcwd()
    path = os.path.join(root_dir, 'files')
    if not os.path.exists(path):
        os.makedirs(path)
    return path

def read_docx(file_path):
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text
    return text

def read_pdf(file_path):
    pdf_file = open(file_path, 'rb')
    pdf_reader = PyPDF2.PdfFileReader(pdf_file)
    text = ''
    for page_num in range(pdf_reader.numPages):
        page = pdf_reader.getPage(page_num)
        text += page.extract_text()
    pdf_file.close()
    return(text)

def read_file(file_path):
    try:
        if(file_path.split('.')[-1]=='pdf'):
            return(read_pdf(file_path))
        if(file_path.split('.')[-1]=='docx'):
            return(read_docx(file_path))
        else:
            raise Exception('Unsupported file extension')
    except:
        raise Exception('Unsupported file')


def summarise_text(giga: GigaChat, text: str) -> str:
    messages = [
        SystemMessage(
            content='Ты научный ассистент. Ты должен  раскрыть основные моменты научной работы по тексту работы и в конце определить практическое применение'
        ),
        HumanMessage(content=f'Раскрой основное содержание статьи по тексту: \"{text}\", а потом отдельно в конце  обязатедьно назови практическое применение')]
    res = giga(messages)
    return res.content

def summarise_topic(giga: GigaChat, topic: str) -> str:
    messages = [
        SystemMessage(
            content='Ты научный ассистент. Ты пишешь научную работу по теме'
        ),
        HumanMessage(content=f'Расскажи возможное содержание работы с названием {topic}. Бери за основу существующие работы с похожей темой. Отдельно вынеси практическое применение работы')]
    res = giga(messages)
    return res.content

def get_profits(giga: GigaChat, text: str) -> str:
    messages = [
        SystemMessage(
            content='Ты научный бизнес-аналитик. Ты должен  прорекламировать потенциальное бизнес применение исследования на основе темы '
        ),
        HumanMessage(content=f'Какое практическое применение и бизнес выгода у исследования: {text}. Так же, оцени экономический эффект от внедрения')]
    res = chat(messages)
    messages.append(res)
    return(res.content)

def make_presention(giga: GigaChat, text: str) -> str:
    messages = [
        SystemMessage(
            content='''Ты помощник по составлению презентаций. 
            Тебе нужно на основе текста составить презентацию в формате JSON.
            В ответе должен быть только валидный JSON и ничего больше. Помни что количество слайдов и текст может изменяться. 
            Вместо полей оформленных как <> тебе нужно подставить текст. 
            Создай 10 слайдов. Если ответ не помещается в сообщение, ограничь его и корректно заверши JSON.
            Вот пример как должен выглядеть ответ на презентацию:
            {
      "title": "<Название презентации>",
      "layout": 1,
      "font": "Calibri",
      "slides": [
        {
          "title": "<Заголовок слайда>",
          "content": "<Текст слайда>",
        },
        {
          "title": "<Заголовок слайда>",
          "content": "<Текст слайда>",
        },
        {
          "title": "<Заголовок слайда>",
          "content": "<Текст слайда>",
        }
      ]
    }
    '''
        )
    ]
    messages.append(HumanMessage(content=text))
    res = chat(messages)
    messages.append(res)
    content = json.loads(res.json())['content']
    print(content)
    path = os.path.join(get_files_dir(), f"{str(uuid.uuid4())}.json")
    with open(path, 'w') as f:
        f.write(content)
    print(f"Файл JSON сохранен {path}")
    return json.loads(content)
    res = chat(messages)
    return res.content

def generate_slides(pres_json):
    # Создаем новую презентацию
    presentation = Presentation()

    for slide_info in pres_json["slides"]:
        # Добавляем слайд с заголовком и содержимым
        slide = presentation.slides.add_slide(presentation.slide_layouts[pres_json['layout']])

        # Добавляем заголовок на слайд
        title = slide.shapes.title
        title.text = slide_info["title"]
        title.font = pres_json['font']
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                # Добавляем содержимое на слайд
        content = slide.placeholders[1]
        content.text = slide_info["content"]
        content.font = pres_json['font']
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
        # Добавление изображений
        image_path = generate_image(slide_info["content"])
        if image_path:
            left = top = Inches(0)
            pic = slide.shapes.add_picture(image_path, left, top, width=presentation.slide_width,
                                           height=presentation.slide_height)
            # This moves it to the background
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            # slide.shapes.add_picture(image_path, 960, 540)
    path = os.path.join(get_files_dir(), f"{str(uuid.uuid4())}.pptx")
    presentation.save(path)

    print(f"Презентация успешно создана и сохранена как {path}")
    return path


